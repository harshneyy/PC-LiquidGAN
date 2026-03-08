from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Physics-Informed LiquidGAN for Zero-Shot Cross-Domain Thermal Image Synthesis"
    subtitle.text = "Harshit Verma (122CS0023)\n\nUnder the Guidance of Dr. N. Srinivas Naik\n\nPrepared for Samsung Innovation Campus (SIC) Review"

    # Slide 2: Problem Statement
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Problem Statement"
    tf = body.text_frame
    tf.text = "Current Limitations in Thermal Image Synthesis:"
    p = tf.add_paragraph()
    p.text = "• Extreme data scarcity in annotated thermal datasets."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Conventional GANs treat thermal maps as arbitrary grayscale images."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Standard models ignore fundamental thermodynamic laws (Heat Diffusion)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Catastrophic mode collapse and spatial degradation in LiquidGAN."
    p.level = 1

    # Slide 3: Proposed Novelty
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Novel Contributions (PC-LiquidGAN)"
    tf = body.text_frame
    tf.text = "Four Key Architectural Innovations:"
    
    p = tf.add_paragraph()
    p.text = "1. Physics-Constrained Loss: Enforces steady-state heat diffusion (∂T/∂t = α∇²T) and energy conservation."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Adaptive ODE Solver (dopri5): Dynamically allocates compute based on thermal gradient complexity (+3.5% SSIM)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. FFT Spectral Loss: Forces generated outputs to obey physical low-frequency dominance in the Fourier spectrum."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. ODE-UNet Generator: Replaced traditional bottleneck latent spaces with spatial ConvODE skip-connections, completely solving mode collapse."
    p.level = 1

    # Slide 4: Results Showcase
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Quantitative Results (ODE-UNet)"
    tf = body.text_frame
    tf.text = "Testing Across 5 Distinct Cross-Disciplinary Domains:"
    
    p = tf.add_paragraph()
    p.text = "• CBSR (Biometric / NIR Face): 0.9976 SSIM | 52.23 dB PSNR"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Agriculture (Chilli Pepper): 0.9947 SSIM | 50.84 dB PSNR"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Agriculture (Tomato Leaf): 0.9945 SSIM | 50.30 dB PSNR"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Medical (Knee X-Ray IR): 0.9631 SSIM | 41.22 dB PSNR"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• KAIST (Surveillance): 0.9351 SSIM | 37.87 dB PSNR"
    p.level = 1

    # Slide 5: Zero Shot Generalization & Demo
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Zero-Shot Generalization & Demo"
    tf = body.text_frame
    tf.text = "Zero-Shot Generalization Test:"
    
    p = tf.add_paragraph()
    p.text = "• The model trained EXCLUSIVELY on Chilli leaves successfully synthesized Tomato leaf thermal maps with ~0.47 SSIM."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Proves the Neural ODE is learning true thermodynamic diffusion paths, not just overfitting source textures."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nInteractive Gradio Dashboard:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• A fully live dashboard was built to allow real-time RGB-to-Thermal synthesis across all 5 models for the defense panel."
    p.level = 1

    prs.save("SIC_Presentation_Harshit_Verma.pptx")
    print("Presentation saved as 'SIC_Presentation_Harshit_Verma.pptx'")

if __name__ == "__main__":
    create_presentation()
